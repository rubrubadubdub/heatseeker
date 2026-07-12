"""Bounded, format-aware document extraction and derived-output storage."""

import hashlib
import io
import json
import zipfile

import pytest
from heatseeker_source_registry.document_processing import (
    DOCM,
    DOCX,
    PDF,
    PPTX,
    XLSX,
    ProcessingLimits,
    detect_media_type,
    process_bytes,
    read_processed_output,
    write_processed_output,
)


def _text_pdf(text: str = "Hello PDF evidence") -> bytes:
    """Construct a tiny, deterministic PDF without a fixture-generation dependency."""
    escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    stream = f"BT /F1 12 Tf 72 720 Td ({escaped}) Tj ET".encode()
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>"
        ),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
    ]
    content = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]
    for number, value in enumerate(objects, start=1):
        offsets.append(len(content))
        content.extend(f"{number} 0 obj\n".encode())
        content.extend(value)
        content.extend(b"\nendobj\n")
    xref_offset = len(content)
    content.extend(f"xref\n0 {len(objects) + 1}\n".encode())
    content.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        content.extend(f"{offset:010d} 00000 n \n".encode())
    content.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode()
    )
    return bytes(content)


def _docx_bytes() -> bytes:
    from docx import Document

    document = Document()
    document.core_properties.title = "Capability statement"
    document.add_heading("Acme Scaffolding", level=1)
    document.add_paragraph("Provides industrial access services.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Service"
    table.cell(1, 0).text = "NSW"
    table.cell(1, 1).text = "Shutdown access"
    output = io.BytesIO()
    document.save(output)
    return output.getvalue()


def _xlsx_bytes() -> bytes:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Companies"
    sheet["A1"] = "Company"
    sheet["B1"] = "Score"
    sheet["A2"] = "Acme"
    sheet["B2"] = "=1+1"
    output = io.BytesIO()
    workbook.save(output)
    workbook.close()
    return output.getvalue()


def _pptx_bytes() -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Project evidence"
    slide.placeholders[1].text = "Ringlock scaffold at the refinery"
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def _png_bytes(width: int = 8, height: int = 6) -> bytes:
    from PIL import Image

    output = io.BytesIO()
    Image.new("RGB", (width, height), "red").save(output, format="PNG")
    return output.getvalue()


def _zip_bytes(entries: dict[str, bytes]) -> bytes:
    output = io.BytesIO()
    with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for name, content in entries.items():
            archive.writestr(name, content)
    return output.getvalue()


@pytest.mark.parametrize(
    ("raw", "declared", "filename", "expected"),
    [
        (_text_pdf(), "text/plain", "wrong.txt", PDF),
        (b"\x89PNG\r\n\x1a\nrest", "application/octet-stream", None, "image/png"),
        (b'{"name":"Acme"}', "application/octet-stream", None, "application/json"),
        (b"name,region\nAcme,NSW\n", None, "companies.csv", "text/csv"),
        (b"plain evidence", None, None, "text/plain"),
        (b"\x00\x01\x02\xff", None, None, "application/octet-stream"),
    ],
)
def test_detect_media_type_prefers_bytes_then_safe_fallbacks(raw, declared, filename, expected):
    assert detect_media_type(raw, declared, filename) == expected


def test_detects_ooxml_package_and_macro_variant():
    assert detect_media_type(_docx_bytes(), filename="statement.bin") == DOCX
    macro = _zip_bytes(
        {
            "word/document.xml": b"<w:document/>",
            "word/vbaProject.bin": b"macro",
        }
    )
    assert detect_media_type(macro, filename="statement.docx") == DOCM
    result = process_bytes(macro, filename="statement.docm")
    assert result.status == "unsupported"
    assert result.error_code == "macro_enabled_document"


def test_html_json_csv_xml_and_plain_text_are_bounded_and_manifested():
    html = process_bytes(
        b"<html><head><title>Acme</title><script>secret()</script></head>"
        b"<body><h1>Projects</h1><p>Refinery scaffold</p></body></html>",
        declared_type="text/html; charset=utf-8",
    )
    assert html.status == "succeeded"
    assert html.text == "# Acme\n\nProjects\nRefinery scaffold"
    assert "secret" not in html.text
    assert html.metadata["title"] == "Acme"
    assert html.manifest["input_sha256"]

    json_result = process_bytes(b'{"companies":["Acme"]}', declared_type="application/json")
    assert json_result.status == "succeeded"
    assert json_result.metadata["json_top_level"] == "dict"

    csv_result = process_bytes(b"name,region\nAcme,NSW\n", filename="evidence.csv")
    assert csv_result.status == "succeeded"
    assert csv_result.metadata["rows_inspected"] == 2

    xml_result = process_bytes(
        b'<?xml version="1.0"?><projects><project>Refinery</project></projects>',
        filename="projects.xml",
    )
    assert xml_result.status == "succeeded"
    assert "Refinery" in xml_result.text

    truncated = process_bytes(b"abcdefghij", limits={"max_text_chars": 5})
    assert truncated.status == "partial"
    assert truncated.text == "abcde"
    assert any("truncated" in warning for warning in truncated.warnings)


def test_malformed_structured_text_returns_corrupt_not_exception():
    result = process_bytes(b'{"broken":', declared_type="application/json")
    assert result.status == "corrupt"
    assert result.error_code == "malformed_text_document"
    assert result.text is None


def test_pdf_extracts_page_text_and_page_locator():
    result = process_bytes(_text_pdf(), filename="capability.pdf")
    assert result.status == "succeeded"
    assert result.detected_content_type == PDF
    assert result.page_count == 1
    assert "Hello PDF evidence" in result.text
    assert result.manifest["segments"][0]["location"]["page"] == 1


def test_encrypted_and_corrupt_pdf_have_explicit_outcomes():
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=100, height=100)
    writer.encrypt("secret")
    output = io.BytesIO()
    writer.write(output)
    encrypted = process_bytes(output.getvalue(), filename="private.pdf")
    assert encrypted.status == "encrypted"
    assert encrypted.error_code == "encrypted_document"

    corrupt = process_bytes(b"%PDF-1.7\nthis is not a PDF", filename="broken.pdf")
    assert corrupt.status == "corrupt"
    assert corrupt.error_code == "parser_error"


def test_docx_extracts_paragraph_table_and_structural_locations():
    result = process_bytes(_docx_bytes(), filename="capability.docx")
    assert result.status == "succeeded"
    assert result.detected_content_type == DOCX
    assert "Provides industrial access services" in result.text
    assert "Shutdown access" in result.text
    kinds = {segment["kind"] for segment in result.manifest["segments"]}
    assert {"paragraph", "table"} <= kinds
    assert result.metadata["title"] == "Capability statement"


def test_xlsx_preserves_formula_without_evaluating_and_has_cell_provenance():
    result = process_bytes(_xlsx_bytes(), filename="companies.xlsx")
    assert result.status == "succeeded"
    assert result.detected_content_type == XLSX
    assert "B2==1+1" in result.text
    assert result.metadata["formulae_preserved_not_evaluated"] == 1
    row = next(segment for segment in result.manifest["segments"] if "B2" in segment["text"])
    assert row["location"]["sheet"] == "Companies"
    assert "B2" in row["location"]["cells"]


def test_pptx_extracts_slide_text_with_slide_and_shape_provenance():
    result = process_bytes(_pptx_bytes(), filename="projects.pptx")
    assert result.status == "succeeded"
    assert result.detected_content_type == PPTX
    assert result.page_count == 1
    assert "Ringlock scaffold at the refinery" in result.text
    assert all(segment["location"]["slide"] == 1 for segment in result.manifest["segments"])


def test_image_metadata_and_pixel_limit():
    raw = _png_bytes(8, 6)
    result = process_bytes(raw, filename="project.png")
    assert result.status == "succeeded"
    assert result.extraction_method == "metadata_only"
    assert result.text is None
    assert result.metadata["width"] == 8
    assert result.metadata["height"] == 6

    blocked = process_bytes(raw, limits={"max_image_pixels": 40})
    assert blocked.status == "quarantined"
    assert blocked.error_code == "image_pixel_limit"


def test_zip_preflight_rejects_traversal_and_excessive_compression():
    traversal = _zip_bytes({"word/document.xml": b"<w:document/>", "../outside.txt": b"bad"})
    result = process_bytes(traversal, filename="unsafe.docx")
    assert result.status == "quarantined"
    assert result.error_code == "unsafe_zip_path"

    compressed = _zip_bytes({"word/document.xml": b"A" * 20_000})
    result = process_bytes(
        compressed,
        filename="bomb.docx",
        limits={"max_zip_compression_ratio": 2.0},
    )
    assert result.status == "quarantined"
    assert result.error_code == "zip_compression_ratio_limit"


def test_input_and_office_limits_return_results_instead_of_raising():
    too_large = process_bytes(b"123456", limits={"max_input_bytes": 5})
    assert too_large.status == "quarantined"
    assert too_large.error_code == "input_size_limit"

    ole = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"legacy"
    legacy = process_bytes(ole, filename="old.doc")
    assert legacy.status == "unsupported"
    assert legacy.error_code == "legacy_office_document"

    invalid_limits = process_bytes(b"hello", limits={"unknown": 1})
    assert invalid_limits.status == "corrupt"
    assert invalid_limits.error_code == "invalid_limits"


def test_versioned_processed_output_is_atomic_readable_and_contained(tmp_path):
    raw = b"raw evidence"
    digest = hashlib.sha256(raw).hexdigest()
    payload = {"status": "succeeded", "text": "Acme"}
    first = write_processed_output(tmp_path, digest, "parser/1", payload)
    second = write_processed_output(tmp_path, digest, "parser/2", payload)

    assert first != second
    assert digest in first
    assert json.loads(read_processed_output(tmp_path, first)) == payload
    assert not list(tmp_path.rglob("*.tmp"))

    with pytest.raises(ValueError, match="escapes"):
        read_processed_output(tmp_path, "../outside")
    with pytest.raises(ValueError, match="simple relative"):
        write_processed_output(tmp_path, digest, "parser/3", b"content", filename="../outside")
    with pytest.raises(ValueError, match="read limit"):
        read_processed_output(tmp_path, first, max_bytes=3)


def test_processing_limits_are_explicit_and_positive():
    limits = ProcessingLimits(max_pages=3)
    limits.validate()
    result = process_bytes(_text_pdf(), limits=limits)
    assert result.manifest["limits"]["max_pages"] == 3
