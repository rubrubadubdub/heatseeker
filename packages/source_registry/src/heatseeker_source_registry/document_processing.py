"""Bounded extraction for untrusted document bytes.

This module deliberately has no database, crawler, or worker dependencies.  Collection
preserves the original bytes first; this module only produces derived text, metadata,
and a provenance manifest.  Optional format libraries are imported lazily so a missing
parser becomes an explicit ``unsupported`` result rather than an import-time failure.

The processor does not execute macros, formulae, external relationships, or embedded
programs.  Legacy Office files and macro-enabled OOXML packages remain preserved raw but
are not interpreted.
"""

from __future__ import annotations

import codecs
import csv
import gzip
import hashlib
import io
import json
import os
import re
import stat
import tempfile
import warnings as stdlib_warnings
import zipfile
from collections.abc import Mapping
from contextlib import suppress
from dataclasses import dataclass, fields
from datetime import date, datetime, time
from html.parser import HTMLParser
from pathlib import Path, PurePosixPath
from typing import Any, ClassVar

PROCESSOR_VERSION = "document-processing/0.1"
MANIFEST_SCHEMA_VERSION = 1

PDF = "application/pdf"
DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
DOCM = "application/vnd.ms-word.document.macroenabled.12"
XLSM = "application/vnd.ms-excel.sheet.macroenabled.12"
PPTM = "application/vnd.ms-powerpoint.presentation.macroenabled.12"
OCTET_STREAM = "application/octet-stream"
ZIP = "application/zip"
OLE_STORAGE = "application/x-ole-storage"

_HEX_64 = re.compile(r"[0-9a-fA-F]{64}\Z")
_CHARSET_RE = re.compile(r"(?:^|;)\s*charset\s*=\s*[\"']?([^;\"'\s]+)", re.I)
_WINDOWS_DRIVE_RE = re.compile(r"^[A-Za-z]:")
_SAFE_OUTPUT_NAME_RE = re.compile(r"[A-Za-z0-9][A-Za-z0-9._-]{0,199}\Z")

_EXTENSION_TYPES = {
    ".txt": "text/plain",
    ".text": "text/plain",
    ".md": "text/markdown",
    ".csv": "text/csv",
    ".json": "application/json",
    ".jsonl": "application/x-ndjson",
    ".xml": "application/xml",
    ".rss": "application/rss+xml",
    ".atom": "application/atom+xml",
    ".html": "text/html",
    ".htm": "text/html",
    ".pdf": PDF,
    ".docx": DOCX,
    ".docm": DOCM,
    ".doc": "application/msword",
    ".xlsx": XLSX,
    ".xlsm": XLSM,
    ".xls": "application/vnd.ms-excel",
    ".pptx": PPTX,
    ".pptm": PPTM,
    ".ppt": "application/vnd.ms-powerpoint",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
    ".tif": "image/tiff",
    ".tiff": "image/tiff",
    ".bmp": "image/bmp",
    ".svg": "image/svg+xml",
    ".zip": ZIP,
}

_DECLARED_ALIASES = {
    "application/x-pdf": PDF,
    "text/xml": "application/xml",
    "application/x-zip-compressed": ZIP,
    "application/vnd.ms-word.document.macroenabled.12": DOCM,
    "application/vnd.ms-excel.sheet.macroenabled.12": XLSM,
    "application/vnd.ms-powerpoint.presentation.macroenabled.12": PPTM,
}

_TEXT_TYPES = {
    "text/plain",
    "text/markdown",
    "text/csv",
    "text/html",
    "application/json",
    "application/x-ndjson",
    "application/xml",
    "application/rss+xml",
    "application/atom+xml",
}

_IMAGE_TYPES = {
    "image/png",
    "image/jpeg",
    "image/gif",
    "image/webp",
    "image/tiff",
    "image/bmp",
}

_MACRO_TYPES = {DOCM, XLSM, PPTM}
_LEGACY_OFFICE_TYPES = {
    OLE_STORAGE,
    "application/msword",
    "application/vnd.ms-excel",
    "application/vnd.ms-powerpoint",
}


@dataclass(frozen=True, slots=True)
class ProcessingLimits:
    """Deterministic resource limits applied before and during extraction."""

    max_input_bytes: int = 10 * 1024 * 1024
    max_text_chars: int = 2_000_000
    max_pages: int = 250
    max_tables: int = 200
    max_table_cells: int = 100_000
    max_spreadsheet_cells: int = 100_000
    max_spreadsheet_rows: int = 10_000
    max_spreadsheet_columns: int = 500
    max_zip_entries: int = 2_000
    max_zip_uncompressed_bytes: int = 100 * 1024 * 1024
    max_zip_entry_bytes: int = 25 * 1024 * 1024
    max_zip_compression_ratio: float = 100.0
    max_attachments: int = 20
    max_image_pixels: int = 40_000_000
    max_image_frames: int = 10
    max_metadata_items: int = 200
    max_error_detail_chars: int = 2_000

    def validate(self) -> None:
        for item in fields(self):
            value = getattr(self, item.name)
            if isinstance(value, bool) or not isinstance(value, (int, float)) or value <= 0:
                raise ValueError(f"{item.name} must be a positive number")


@dataclass(frozen=True, slots=True)
class ProcessingResult:
    status: str
    detected_content_type: str
    extraction_method: str
    text: str | None
    page_count: int | None
    metadata: dict[str, Any]
    warnings: list[str]
    error_code: str | None
    error_detail: str | None
    manifest: dict[str, Any]


@dataclass(slots=True)
class _ParserOutput:
    status: str = "succeeded"
    extraction_method: str = "native"
    text: str | None = None
    page_count: int | None = None
    metadata: dict[str, Any] | None = None
    warnings: list[str] | None = None
    error_code: str | None = None
    error_detail: str | None = None
    segments: list[dict[str, Any]] | None = None


class _DependencyUnavailable(RuntimeError):
    pass


class _ZipSafetyError(RuntimeError):
    def __init__(self, code: str, detail: str, *, status: str = "quarantined") -> None:
        super().__init__(detail)
        self.code = code
        self.detail = detail
        self.status = status


class _OutputBuilder:
    def __init__(self, limit: int) -> None:
        self.limit = limit
        self.length = 0
        self.parts: list[str] = []
        self.truncated = False

    def append(self, value: str) -> str:
        if not value:
            return ""
        remaining = self.limit - self.length
        if remaining <= 0:
            self.truncated = True
            return ""
        added = value[:remaining]
        self.parts.append(added)
        self.length += len(added)
        if len(added) < len(value):
            self.truncated = True
        return added

    def text(self) -> str:
        return "".join(self.parts).strip()


def _normalise_declared_type(declared_type: str | None) -> str | None:
    if not declared_type:
        return None
    media_type = declared_type.split(";", 1)[0].strip().lower()
    if not media_type:
        return None
    return _DECLARED_ALIASES.get(media_type, media_type)


def _filename_type(filename: str | None) -> str | None:
    if not filename:
        return None
    suffix = Path(filename.replace("\\", "/")).suffix.lower()
    return _EXTENSION_TYPES.get(suffix)


def _zip_package_type(raw: bytes) -> str | None:
    try:
        with zipfile.ZipFile(io.BytesIO(raw)) as archive:
            names = {item.filename.replace("\\", "/").lower() for item in archive.infolist()}
    except (OSError, ValueError, zipfile.BadZipFile):
        return None
    has_macros = any(name.endswith("vbaproject.bin") for name in names)
    if "word/document.xml" in names:
        return DOCM if has_macros else DOCX
    if "xl/workbook.xml" in names:
        return XLSM if has_macros else XLSX
    if "ppt/presentation.xml" in names:
        return PPTM if has_macros else PPTX
    return ZIP


def _looks_like_text(raw: bytes) -> bool:
    if not raw:
        return True
    sample = raw[:8192]
    if b"\x00" in sample and not sample.startswith(
        (codecs.BOM_UTF16_LE, codecs.BOM_UTF16_BE, codecs.BOM_UTF32_LE, codecs.BOM_UTF32_BE)
    ):
        return False
    try:
        sample.decode("utf-8")
    except UnicodeDecodeError:
        printable = sum(byte in b"\t\n\r" or 32 <= byte < 127 for byte in sample)
        return printable / len(sample) >= 0.85
    return True


def detect_media_type(
    raw: bytes, declared_type: str | None = None, filename: str | None = None
) -> str:
    """Detect a conservative media type, preferring byte signatures over metadata.

    The function never opens or expands archive members.  OOXML identification reads
    only the ZIP central directory and well-known member names.
    """

    try:
        content = bytes(raw)
    except (TypeError, ValueError):
        return OCTET_STREAM

    head = content[:1024]
    stripped = head.lstrip()
    lowered = stripped.lower()

    if b"%pdf-" in lowered[:128]:
        return PDF
    if content.startswith(b"\x89PNG\r\n\x1a\n"):
        return "image/png"
    if content.startswith(b"\xff\xd8\xff"):
        return "image/jpeg"
    if content.startswith((b"GIF87a", b"GIF89a")):
        return "image/gif"
    if len(content) >= 12 and content.startswith(b"RIFF") and content[8:12] == b"WEBP":
        return "image/webp"
    if content.startswith((b"II*\x00", b"MM\x00*")):
        return "image/tiff"
    if content.startswith(b"BM"):
        return "image/bmp"
    if content.startswith(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"):
        return _filename_type(filename) or OLE_STORAGE
    if content.startswith((b"PK\x03\x04", b"PK\x05\x06", b"PK\x07\x08")):
        return _zip_package_type(content) or _filename_type(filename) or ZIP

    if lowered.startswith((b"<!doctype html", b"<html", b"<head", b"<body")):
        return "text/html"
    if lowered.startswith(b"<?xml"):
        return _filename_type(filename) or "application/xml"
    if lowered.startswith((b"<rss", b"<feed")):
        return "application/xml"
    if lowered.startswith((b"{", b"[")) and _looks_like_text(content):
        return "application/json"

    declared = _normalise_declared_type(declared_type)
    if declared and declared != OCTET_STREAM:
        return declared
    from_filename = _filename_type(filename)
    if from_filename:
        return from_filename
    if _looks_like_text(content):
        return "text/plain"
    return declared or OCTET_STREAM


def _coerce_limits(limits: ProcessingLimits | Mapping[str, Any] | None) -> ProcessingLimits:
    if limits is None:
        result = ProcessingLimits()
    elif isinstance(limits, ProcessingLimits):
        result = limits
    elif isinstance(limits, Mapping):
        allowed = {item.name for item in fields(ProcessingLimits)}
        unknown = sorted(set(limits) - allowed)
        if unknown:
            raise ValueError(f"unknown processing limits: {', '.join(unknown)}")
        result = ProcessingLimits(**dict(limits))
    else:
        raise TypeError("limits must be ProcessingLimits, a mapping, or None")
    result.validate()
    return result


def _display_filename(filename: str | None) -> str | None:
    if not filename:
        return None
    value = "".join(char for char in filename if char >= " " and char != "\x7f").strip()
    return value[:500] or None


def _safe_scalar(value: Any, *, max_chars: int = 2_000) -> Any:
    if value is None or isinstance(value, (bool, int, float)):
        return value
    if isinstance(value, (datetime, date, time)):
        return value.isoformat()
    if isinstance(value, bytes):
        return f"<bytes:{len(value)}>"
    text = str(value)
    return text[:max_chars]


def _safe_metadata(metadata: Mapping[str, Any] | None, limit: int) -> dict[str, Any]:
    result: dict[str, Any] = {}
    for index, (key, value) in enumerate((metadata or {}).items()):
        if index >= limit:
            result["_truncated"] = True
            break
        safe_key = str(key)[:200]
        if isinstance(value, Mapping):
            result[safe_key] = {
                str(k)[:200]: _safe_scalar(v) for k, v in list(value.items())[:limit]
            }
        elif isinstance(value, (list, tuple, set)):
            result[safe_key] = [_safe_scalar(item) for item in list(value)[:limit]]
        else:
            result[safe_key] = _safe_scalar(value)
    return result


def _bounded_detail(value: str | None, limits: ProcessingLimits) -> str | None:
    if value is None:
        return None
    return value.replace("\x00", "")[: limits.max_error_detail_chars]


def _exception_detail(exc: Exception, limits: ProcessingLimits) -> str:
    message = str(exc).replace("\x00", "")
    return f"{type(exc).__name__}: {message}"[: limits.max_error_detail_chars]


def _to_result(
    raw: bytes,
    detected: str,
    declared_type: str | None,
    filename: str | None,
    output: _ParserOutput,
    limits: ProcessingLimits,
) -> ProcessingResult:
    metadata = _safe_metadata(output.metadata, limits.max_metadata_items)
    display_name = _display_filename(filename)
    if display_name:
        metadata.setdefault("original_filename", display_name)
    warning_values = [str(item)[:1_000] for item in (output.warnings or [])]
    error_detail = _bounded_detail(output.error_detail, limits)
    segments = output.segments or []
    manifest: dict[str, Any] = {
        "schema_version": MANIFEST_SCHEMA_VERSION,
        "processor_version": PROCESSOR_VERSION,
        "input_sha256": hashlib.sha256(raw).hexdigest(),
        "input_bytes": len(raw),
        "declared_content_type": declared_type,
        "detected_content_type": detected,
        "status": output.status,
        "extraction_method": output.extraction_method,
        "page_count": output.page_count,
        "metadata": metadata,
        "warnings": warning_values,
        "error_code": output.error_code,
        "error_detail": error_detail,
        "segments": segments,
        "limits": {
            "max_input_bytes": limits.max_input_bytes,
            "max_text_chars": limits.max_text_chars,
            "max_pages": limits.max_pages,
        },
    }
    return ProcessingResult(
        status=output.status,
        detected_content_type=detected,
        extraction_method=output.extraction_method,
        text=output.text,
        page_count=output.page_count,
        metadata=metadata,
        warnings=warning_values,
        error_code=output.error_code,
        error_detail=error_detail,
        manifest=manifest,
    )


def _failure(
    status: str,
    code: str,
    detail: str,
    *,
    warnings: list[str] | None = None,
    metadata: dict[str, Any] | None = None,
) -> _ParserOutput:
    return _ParserOutput(
        status=status,
        extraction_method="none",
        metadata=metadata or {},
        warnings=warnings or [],
        error_code=code,
        error_detail=detail,
        segments=[],
    )


def _zip_preflight(raw: bytes, limits: ProcessingLimits) -> dict[str, Any]:
    try:
        archive = zipfile.ZipFile(io.BytesIO(raw))
    except (OSError, ValueError, zipfile.BadZipFile) as exc:
        raise _ZipSafetyError("corrupt_zip", str(exc), status="corrupt") from exc

    with archive:
        entries = archive.infolist()
        if len(entries) > limits.max_zip_entries:
            raise _ZipSafetyError(
                "zip_entry_limit",
                f"archive contains {len(entries)} entries; limit is {limits.max_zip_entries}",
            )
        total = 0
        for entry in entries:
            name = entry.filename.replace("\\", "/")
            path = PurePosixPath(name)
            if (
                not name
                or name.startswith("/")
                or path.is_absolute()
                or ".." in path.parts
                or _WINDOWS_DRIVE_RE.match(name)
            ):
                raise _ZipSafetyError("unsafe_zip_path", f"unsafe archive member path: {name!r}")
            unix_mode = (entry.external_attr >> 16) & 0xFFFF
            if unix_mode and stat.S_ISLNK(unix_mode):
                raise _ZipSafetyError(
                    "zip_symlink", f"symbolic link member is not allowed: {name!r}"
                )
            if entry.flag_bits & 0x1:
                raise _ZipSafetyError(
                    "encrypted_archive",
                    f"encrypted archive member is not supported: {name!r}",
                    status="encrypted",
                )
            if entry.file_size > limits.max_zip_entry_bytes:
                raise _ZipSafetyError(
                    "zip_entry_size_limit",
                    f"archive member {name!r} declares {entry.file_size} bytes; "
                    f"limit is {limits.max_zip_entry_bytes}",
                )
            total += entry.file_size
            if total > limits.max_zip_uncompressed_bytes:
                raise _ZipSafetyError(
                    "zip_total_size_limit",
                    f"archive declares more than {limits.max_zip_uncompressed_bytes} "
                    "uncompressed bytes",
                )
            if entry.file_size >= 1024:
                ratio = entry.file_size / max(entry.compress_size, 1)
                if ratio > limits.max_zip_compression_ratio:
                    raise _ZipSafetyError(
                        "zip_compression_ratio_limit",
                        f"archive member {name!r} has compression ratio {ratio:.1f}; "
                        f"limit is {limits.max_zip_compression_ratio:.1f}",
                    )
        return {"zip_entries": len(entries), "zip_uncompressed_bytes": total}


def _declared_charset(declared_type: str | None) -> str | None:
    if not declared_type:
        return None
    match = _CHARSET_RE.search(declared_type)
    return match.group(1) if match else None


def _decode_text(raw: bytes, declared_type: str | None) -> tuple[str, str, list[str]]:
    warnings: list[str] = []
    candidates: list[str] = []
    if raw.startswith(codecs.BOM_UTF8):
        candidates.append("utf-8-sig")
    elif raw.startswith((codecs.BOM_UTF32_LE, codecs.BOM_UTF32_BE)):
        candidates.append("utf-32")
    elif raw.startswith((codecs.BOM_UTF16_LE, codecs.BOM_UTF16_BE)):
        candidates.append("utf-16")
    declared_charset = _declared_charset(declared_type)
    if declared_charset:
        try:
            codecs.lookup(declared_charset)
        except LookupError:
            warnings.append(f"unknown declared charset {declared_charset!r}; ignored")
        else:
            candidates.append(declared_charset)
    candidates.append("utf-8")
    tried: set[str] = set()
    for encoding in candidates:
        normalised = encoding.lower()
        if normalised in tried:
            continue
        tried.add(normalised)
        try:
            return raw.decode(encoding), encoding, warnings
        except UnicodeDecodeError:
            continue
    warnings.append("text was not UTF-8; decoded as Windows-1252")
    return raw.decode("cp1252", errors="replace"), "cp1252", warnings


class _HTMLExtractor(HTMLParser):
    _DROP: ClassVar[set[str]] = {
        "script",
        "style",
        "noscript",
        "svg",
        "iframe",
        "canvas",
        "template",
    }
    _BLOCK: ClassVar[set[str]] = {
        "address",
        "article",
        "aside",
        "blockquote",
        "br",
        "div",
        "dl",
        "dt",
        "dd",
        "figcaption",
        "figure",
        "h1",
        "h2",
        "h3",
        "h4",
        "h5",
        "h6",
        "li",
        "main",
        "ol",
        "p",
        "pre",
        "section",
        "table",
        "td",
        "th",
        "tr",
        "ul",
    }

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.drop_depth = 0
        self.in_title = False
        self.title_parts: list[str] = []
        self.parts: list[str] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        del attrs
        lowered = tag.lower()
        if self.drop_depth:
            if lowered in self._DROP:
                self.drop_depth += 1
            return
        if lowered in self._DROP:
            self.drop_depth = 1
            return
        if lowered == "title":
            self.in_title = True
            return
        if lowered in self._BLOCK:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        lowered = tag.lower()
        if self.drop_depth:
            if lowered in self._DROP:
                self.drop_depth -= 1
            return
        if lowered == "title":
            self.in_title = False
            return
        if lowered in self._BLOCK:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if self.drop_depth:
            return
        if self.in_title:
            self.title_parts.append(data)
        else:
            self.parts.append(data)

    def result(self) -> tuple[str, str | None]:
        title = " ".join(" ".join(self.title_parts).split()) or None
        text = "".join(self.parts)
        lines = [" ".join(line.split()) for line in text.splitlines()]
        text = "\n".join(line for line in lines if line).strip()
        if title:
            text = f"# {title}\n\n{text}" if text else f"# {title}"
        return text, title


def _extract_html(text: str) -> tuple[str, str | None]:
    parser = _HTMLExtractor()
    parser.feed(text)
    parser.close()
    return parser.result()


def _extract_xml(text: str) -> tuple[str, dict[str, Any], list[str]]:
    warnings: list[str] = []
    try:
        from defusedxml import ElementTree as safe_element_tree
    except ImportError:
        warnings.append("defusedxml unavailable; conservative markup stripping used")
        extracted, _ = _extract_html(text)
        return extracted, {}, warnings
    root = safe_element_tree.fromstring(text)
    values: list[str] = []
    for element in root.iter():
        if element.text and element.text.strip():
            values.append(" ".join(element.text.split()))
        if element.tail and element.tail.strip():
            values.append(" ".join(element.tail.split()))
    return "\n".join(values), {"root_element": str(root.tag)[:500]}, warnings


def _process_textual(
    raw: bytes, detected: str, declared_type: str | None, limits: ProcessingLimits
) -> _ParserOutput:
    decoded, encoding, decode_warnings = _decode_text(raw, declared_type)
    decoded = decoded.replace("\r\n", "\n").replace("\r", "\n")
    metadata: dict[str, Any] = {"encoding": encoding, "source_characters": len(decoded)}
    warnings = list(decode_warnings)

    try:
        if detected == "text/html":
            extracted, title = _extract_html(decoded)
            if title:
                metadata["title"] = title
        elif detected in {"application/xml", "application/rss+xml", "application/atom+xml"}:
            extracted, xml_metadata, xml_warnings = _extract_xml(decoded)
            metadata.update(xml_metadata)
            warnings.extend(xml_warnings)
        elif detected in {"application/json", "application/x-ndjson"}:
            if detected == "application/json":
                parsed = json.loads(decoded)
                metadata["json_top_level"] = type(parsed).__name__
                if isinstance(parsed, (dict, list)):
                    metadata["json_top_level_items"] = len(parsed)
            else:
                lines = [line for line in decoded.splitlines() if line.strip()]
                for line in lines[: limits.max_spreadsheet_rows]:
                    json.loads(line)
                metadata["json_lines"] = len(lines)
                if len(lines) > limits.max_spreadsheet_rows:
                    warnings.append("JSON Lines validation stopped at the configured row limit")
            extracted = decoded.strip()
        elif detected == "text/csv":
            rows = 0
            cells = 0
            max_columns = 0
            for row in csv.reader(io.StringIO(decoded)):
                rows += 1
                cells += len(row)
                max_columns = max(max_columns, len(row))
                if rows >= limits.max_spreadsheet_rows or cells >= limits.max_spreadsheet_cells:
                    warnings.append("CSV inspection stopped at the configured row/cell limit")
                    break
            metadata.update(
                {
                    "rows_inspected": rows,
                    "cells_inspected": cells,
                    "max_columns": max_columns,
                }
            )
            extracted = decoded.strip()
        else:
            extracted = decoded.strip()
    except Exception as exc:
        return _failure("corrupt", "malformed_text_document", _exception_detail(exc, limits))

    builder = _OutputBuilder(limits.max_text_chars)
    kept = builder.append(extracted)
    if builder.truncated:
        warnings.append(f"text truncated at {limits.max_text_chars} characters")
    status = "partial" if builder.truncated else "succeeded"
    return _ParserOutput(
        status=status,
        extraction_method="native",
        text=builder.text(),
        metadata=metadata,
        warnings=warnings,
        segments=[
            {
                "id": "document:text:1",
                "kind": "text",
                "method": "native",
                "location": {"scope": "document"},
                "text": kept,
            }
        ]
        if kept
        else [],
    )


def _pdf_metadata(reader: Any, limits: ProcessingLimits) -> dict[str, Any]:
    result: dict[str, Any] = {}
    try:
        values = reader.metadata or {}
        for index, (key, value) in enumerate(values.items()):
            if index >= limits.max_metadata_items:
                result["metadata_truncated"] = True
                break
            result[str(key).lstrip("/")[:200]] = _safe_scalar(value)
    except Exception as exc:
        result["metadata_error"] = type(exc).__name__
    return result


def _process_pdf(raw: bytes, limits: ProcessingLimits) -> _ParserOutput:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise _DependencyUnavailable("pypdf is not installed") from exc

    reader = PdfReader(io.BytesIO(raw), strict=False)
    if reader.is_encrypted:
        return _failure(
            "encrypted",
            "encrypted_document",
            "encrypted PDFs require an explicit password workflow and were not interpreted",
        )

    page_count = len(reader.pages)
    metadata = _pdf_metadata(reader, limits)
    attachment_names: list[str] = []
    try:
        attachment_list = getattr(reader, "attachment_list", None)
        if attachment_list is not None:
            for index, attachment in enumerate(attachment_list):
                if index >= limits.max_attachments:
                    break
                attachment_names.append(str(getattr(attachment, "name", "attachment"))[:500])
    except Exception:
        attachment_names = []
    if attachment_names:
        metadata["attachment_names"] = attachment_names

    builder = _OutputBuilder(limits.max_text_chars)
    segments: list[dict[str, Any]] = []
    warnings: list[str] = []
    failed_pages = 0
    pages_to_process = min(page_count, limits.max_pages)
    if page_count > limits.max_pages:
        warnings.append(f"only the first {limits.max_pages} of {page_count} pages were processed")
    for page_index in range(pages_to_process):
        try:
            page_text = reader.pages[page_index].extract_text() or ""
        except Exception as exc:
            failed_pages += 1
            warnings.append(
                f"page {page_index + 1} native text extraction failed: {type(exc).__name__}"
            )
            continue
        page_text = page_text.replace("\r\n", "\n").replace("\r", "\n").strip()
        if not page_text:
            continue
        builder.append(f"\n\n--- page {page_index + 1} ---\n\n")
        kept = builder.append(page_text)
        if kept:
            segments.append(
                {
                    "id": f"page:{page_index + 1}:text:1",
                    "kind": "text",
                    "method": "native",
                    "location": {"page": page_index + 1, "coordinate_space": "pdf_points"},
                    "text": kept,
                }
            )
        if builder.truncated:
            break
    if builder.truncated:
        warnings.append(f"text truncated at {limits.max_text_chars} characters")
    if not segments:
        warnings.append("no extractable native text; OCR was not performed")
    partial = page_count > limits.max_pages or builder.truncated or failed_pages > 0 or not segments
    return _ParserOutput(
        status="partial" if partial else "succeeded",
        extraction_method="native",
        text=builder.text(),
        page_count=page_count,
        metadata=metadata,
        warnings=warnings,
        segments=segments,
    )


def _office_properties(properties: Any) -> dict[str, Any]:
    result: dict[str, Any] = {}
    for name in (
        "title",
        "subject",
        "author",
        "keywords",
        "comments",
        "category",
        "created",
        "modified",
        "last_modified_by",
        "revision",
    ):
        try:
            value = getattr(properties, name)
        except (AttributeError, ValueError):
            continue
        if value not in (None, ""):
            result[name] = _safe_scalar(value)
    return result


def _add_segment(
    builder: _OutputBuilder,
    segments: list[dict[str, Any]],
    text: str,
    *,
    segment_id: str,
    kind: str,
    location: dict[str, Any],
    marker: str = "",
) -> bool:
    cleaned = text.replace("\r\n", "\n").replace("\r", "\n").strip()
    if not cleaned:
        return True
    if marker:
        builder.append(marker)
    kept = builder.append(cleaned)
    builder.append("\n")
    if kept:
        segments.append(
            {
                "id": segment_id,
                "kind": kind,
                "method": "native",
                "location": location,
                "text": kept,
            }
        )
    return not builder.truncated


def _docx_table_text(
    table: Any, limits: ProcessingLimits, cells_seen: int
) -> tuple[str, int, bool]:
    lines: list[str] = []
    truncated = False
    for row in table.rows:
        values: list[str] = []
        for cell in row.cells:
            if cells_seen >= limits.max_table_cells:
                truncated = True
                break
            cells_seen += 1
            values.append(" ".join(cell.text.split()))
        if values:
            lines.append("\t".join(values))
        if truncated:
            break
    return "\n".join(lines), cells_seen, truncated


def _process_docx(raw: bytes, limits: ProcessingLimits) -> _ParserOutput:
    try:
        from docx import Document
        from docx.table import Table
    except ImportError as exc:
        raise _DependencyUnavailable("python-docx is not installed") from exc

    document = Document(io.BytesIO(raw))
    builder = _OutputBuilder(limits.max_text_chars)
    segments: list[dict[str, Any]] = []
    warnings: list[str] = []
    tables_seen = 0
    cells_seen = 0
    paragraph_index = 0

    blocks = (
        document.iter_inner_content()
        if hasattr(document, "iter_inner_content")
        else [*document.paragraphs, *document.tables]
    )
    for block in blocks:
        if isinstance(block, Table):
            if tables_seen >= limits.max_tables:
                warnings.append("DOCX table extraction stopped at the configured table limit")
                break
            tables_seen += 1
            table_text, cells_seen, table_truncated = _docx_table_text(block, limits, cells_seen)
            if not _add_segment(
                builder,
                segments,
                table_text,
                segment_id=f"body:table:{tables_seen}",
                kind="table",
                location={"part": "body", "table": tables_seen},
                marker=f"\n--- table {tables_seen} ---\n",
            ):
                break
            if table_truncated:
                warnings.append("DOCX table extraction stopped at the configured cell limit")
                break
        else:
            paragraph_index += 1
            style = None
            with suppress(AttributeError, ValueError):
                style = block.style.name
            if not _add_segment(
                builder,
                segments,
                block.text,
                segment_id=f"body:paragraph:{paragraph_index}",
                kind="paragraph",
                location={"part": "body", "paragraph": paragraph_index, "style": style},
            ):
                break

    for section_index, section in enumerate(document.sections, start=1):
        for part_name, part in (("header", section.header), ("footer", section.footer)):
            for index, paragraph in enumerate(part.paragraphs, start=1):
                if not _add_segment(
                    builder,
                    segments,
                    paragraph.text,
                    segment_id=f"section:{section_index}:{part_name}:paragraph:{index}",
                    kind="paragraph",
                    location={
                        "part": part_name,
                        "section": section_index,
                        "paragraph": index,
                    },
                    marker=f"\n--- {part_name} section {section_index} ---\n",
                ):
                    break
    if builder.truncated:
        warnings.append(f"text truncated at {limits.max_text_chars} characters")
    metadata = _office_properties(document.core_properties)
    metadata.update(
        {
            "paragraphs_inspected": paragraph_index,
            "tables_inspected": tables_seen,
            "table_cells_inspected": cells_seen,
            "inline_images": len(document.inline_shapes),
        }
    )
    return _ParserOutput(
        status="partial" if warnings or builder.truncated else "succeeded",
        extraction_method="native",
        text=builder.text(),
        metadata=metadata,
        warnings=warnings,
        segments=segments,
    )


def _cell_text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, (datetime, date, time)):
        return value.isoformat()
    return str(value)


def _process_xlsx(raw: bytes, limits: ProcessingLimits) -> _ParserOutput:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise _DependencyUnavailable("openpyxl is not installed") from exc

    workbook = load_workbook(
        io.BytesIO(raw),
        read_only=True,
        data_only=False,
        keep_links=False,
        keep_vba=False,
    )
    builder = _OutputBuilder(limits.max_text_chars)
    segments: list[dict[str, Any]] = []
    warnings: list[str] = []
    cells_seen = 0
    formulae_seen = 0
    rows_seen = 0
    stopped = False
    try:
        for sheet_index, worksheet in enumerate(workbook.worksheets, start=1):
            builder.append(f"\n--- sheet: {worksheet.title} ---\n")
            for row_index, row in enumerate(worksheet.iter_rows(), start=1):
                if row_index > limits.max_spreadsheet_rows:
                    warnings.append(
                        f"sheet {worksheet.title!r} stopped at row {limits.max_spreadsheet_rows}"
                    )
                    stopped = True
                    break
                row_values: list[str] = []
                coordinates: list[str] = []
                for column_index, cell in enumerate(row, start=1):
                    if column_index > limits.max_spreadsheet_columns:
                        warnings.append(
                            f"sheet {worksheet.title!r} row {row_index} stopped at column "
                            f"{limits.max_spreadsheet_columns}"
                        )
                        stopped = True
                        break
                    if cell.value is None:
                        continue
                    if cells_seen >= limits.max_spreadsheet_cells:
                        warnings.append("XLSX extraction stopped at the configured cell limit")
                        stopped = True
                        break
                    cells_seen += 1
                    value = _cell_text(cell.value)
                    if value.startswith("="):
                        formulae_seen += 1
                    coordinates.append(cell.coordinate)
                    row_values.append(f"{cell.coordinate}={value}")
                if row_values:
                    rows_seen += 1
                    if not _add_segment(
                        builder,
                        segments,
                        "\t".join(row_values),
                        segment_id=f"sheet:{sheet_index}:row:{row_index}",
                        kind="spreadsheet_row",
                        location={
                            "sheet": worksheet.title,
                            "sheet_index": sheet_index,
                            "row": row_index,
                            "cells": coordinates,
                        },
                    ):
                        stopped = True
                if stopped:
                    break
            if stopped:
                break
    finally:
        workbook.close()
    if builder.truncated:
        warnings.append(f"text truncated at {limits.max_text_chars} characters")
    metadata = _office_properties(workbook.properties)
    metadata.update(
        {
            "sheet_count": len(workbook.sheetnames),
            "rows_inspected": rows_seen,
            "nonempty_cells_inspected": cells_seen,
            "formulae_preserved_not_evaluated": formulae_seen,
        }
    )
    return _ParserOutput(
        status="partial" if warnings or stopped or builder.truncated else "succeeded",
        extraction_method="native",
        text=builder.text(),
        metadata=metadata,
        warnings=warnings,
        segments=segments,
    )


def _shape_location(shape: Any, slide_index: int, shape_index: int) -> dict[str, Any]:
    location: dict[str, Any] = {"slide": slide_index, "shape": shape_index}
    for name in ("left", "top", "width", "height"):
        with suppress(AttributeError, TypeError, ValueError):
            location[f"{name}_emu"] = int(getattr(shape, name))
    return location


def _process_pptx(raw: bytes, limits: ProcessingLimits) -> _ParserOutput:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise _DependencyUnavailable("python-pptx is not installed") from exc

    presentation = Presentation(io.BytesIO(raw))
    builder = _OutputBuilder(limits.max_text_chars)
    segments: list[dict[str, Any]] = []
    warnings: list[str] = []
    table_count = 0
    table_cells = 0
    slides_to_process = min(len(presentation.slides), limits.max_pages)
    if len(presentation.slides) > limits.max_pages:
        warnings.append(
            f"only the first {limits.max_pages} of {len(presentation.slides)} slides were processed"
        )
    for slide_index in range(1, slides_to_process + 1):
        slide = presentation.slides[slide_index - 1]
        builder.append(f"\n--- slide {slide_index} ---\n")
        for shape_index, shape in enumerate(slide.shapes, start=1):
            location = _shape_location(shape, slide_index, shape_index)
            if getattr(shape, "has_text_frame", False) and not _add_segment(
                builder,
                segments,
                shape.text,
                segment_id=f"slide:{slide_index}:shape:{shape_index}:text",
                kind="slide_text",
                location=location,
            ):
                break
            if getattr(shape, "has_table", False):
                if table_count >= limits.max_tables:
                    warnings.append("PPTX table extraction stopped at the configured table limit")
                    break
                table_count += 1
                lines: list[str] = []
                stopped = False
                for row in shape.table.rows:
                    values: list[str] = []
                    for cell in row.cells:
                        if table_cells >= limits.max_table_cells:
                            stopped = True
                            break
                        table_cells += 1
                        values.append(" ".join(cell.text.split()))
                    if values:
                        lines.append("\t".join(values))
                    if stopped:
                        break
                if not _add_segment(
                    builder,
                    segments,
                    "\n".join(lines),
                    segment_id=f"slide:{slide_index}:shape:{shape_index}:table",
                    kind="table",
                    location=location,
                ):
                    break
                if stopped:
                    warnings.append("PPTX table extraction stopped at the configured cell limit")
                    break
        if builder.truncated:
            break
        try:
            has_notes = slide.has_notes_slide
        except AttributeError:
            has_notes = False
        if has_notes:
            try:
                notes = slide.notes_slide.notes_text_frame.text
            except (AttributeError, ValueError):
                notes = ""
            _add_segment(
                builder,
                segments,
                notes,
                segment_id=f"slide:{slide_index}:notes",
                kind="notes",
                location={"slide": slide_index, "part": "notes"},
                marker="\n--- notes ---\n",
            )
    if builder.truncated:
        warnings.append(f"text truncated at {limits.max_text_chars} characters")
    metadata = _office_properties(presentation.core_properties)
    metadata.update(
        {
            "slide_count": len(presentation.slides),
            "tables_inspected": table_count,
            "table_cells_inspected": table_cells,
        }
    )
    return _ParserOutput(
        status="partial" if warnings or builder.truncated else "succeeded",
        extraction_method="native",
        text=builder.text(),
        page_count=len(presentation.slides),
        metadata=metadata,
        warnings=warnings,
        segments=segments,
    )


def _process_image(raw: bytes, limits: ProcessingLimits) -> _ParserOutput:
    try:
        from PIL import ExifTags, Image
    except ImportError as exc:
        raise _DependencyUnavailable("Pillow is not installed") from exc

    warning_values: list[str] = []
    with stdlib_warnings.catch_warnings():
        stdlib_warnings.simplefilter("error", Image.DecompressionBombWarning)
        with Image.open(io.BytesIO(raw)) as image:
            width, height = image.size
            pixels = width * height
            frame_count = int(getattr(image, "n_frames", 1))
            image_format = image.format
            mode = image.mode
            if pixels > limits.max_image_pixels:
                return _failure(
                    "quarantined",
                    "image_pixel_limit",
                    f"image contains {pixels} pixels; limit is {limits.max_image_pixels}",
                    metadata={"width": width, "height": height, "pixels": pixels},
                )
            image.verify()
    if frame_count > limits.max_image_frames:
        warning_values.append(
            f"image has {frame_count} frames; only metadata was read and frame limit is "
            f"{limits.max_image_frames}"
        )

    exif: dict[str, Any] = {}
    with stdlib_warnings.catch_warnings():
        stdlib_warnings.simplefilter("error", Image.DecompressionBombWarning)
        with Image.open(io.BytesIO(raw)) as image:
            try:
                values = image.getexif()
                for index, (key, value) in enumerate(values.items()):
                    if index >= limits.max_metadata_items:
                        exif["_truncated"] = True
                        break
                    name = ExifTags.TAGS.get(key, str(key))
                    exif[str(name)[:200]] = _safe_scalar(value)
            except (AttributeError, OSError, ValueError):
                pass
    metadata: dict[str, Any] = {
        "format": image_format,
        "width": width,
        "height": height,
        "pixels": pixels,
        "mode": mode,
        "frame_count": frame_count,
    }
    if exif:
        metadata["exif"] = exif
    return _ParserOutput(
        status="partial" if warning_values else "succeeded",
        extraction_method="metadata_only",
        page_count=frame_count if frame_count > 1 else None,
        metadata=metadata,
        warnings=warning_values,
        segments=[],
    )


def process_bytes(
    raw: bytes,
    declared_type: str | None = None,
    filename: str | None = None,
    limits: ProcessingLimits | Mapping[str, Any] | None = None,
) -> ProcessingResult:
    """Process untrusted bytes into bounded derived data without raising parser errors."""

    try:
        content = bytes(raw)
    except (TypeError, ValueError) as exc:
        fallback_limits = ProcessingLimits()
        return _to_result(
            b"",
            OCTET_STREAM,
            declared_type,
            filename,
            _failure("corrupt", "invalid_input", _exception_detail(exc, fallback_limits)),
            fallback_limits,
        )

    detected = detect_media_type(content, declared_type, filename)
    try:
        active_limits = _coerce_limits(limits)
    except (TypeError, ValueError) as exc:
        fallback_limits = ProcessingLimits()
        return _to_result(
            content,
            detected,
            declared_type,
            filename,
            _failure("corrupt", "invalid_limits", _exception_detail(exc, fallback_limits)),
            fallback_limits,
        )

    if len(content) > active_limits.max_input_bytes:
        output = _failure(
            "quarantined",
            "input_size_limit",
            f"input contains {len(content)} bytes; limit is {active_limits.max_input_bytes}",
        )
        return _to_result(content, detected, declared_type, filename, output, active_limits)

    zip_metadata: dict[str, Any] = {}
    if detected in {ZIP, DOCX, XLSX, PPTX, DOCM, XLSM, PPTM}:
        try:
            zip_metadata = _zip_preflight(content, active_limits)
        except _ZipSafetyError as exc:
            output = _failure(exc.status, exc.code, exc.detail)
            return _to_result(content, detected, declared_type, filename, output, active_limits)

    if detected in _MACRO_TYPES:
        output = _failure(
            "unsupported",
            "macro_enabled_document",
            "macro-enabled Office documents are preserved raw but are not interpreted",
            metadata=zip_metadata,
        )
        return _to_result(content, detected, declared_type, filename, output, active_limits)
    if detected in _LEGACY_OFFICE_TYPES:
        output = _failure(
            "unsupported",
            "legacy_office_document",
            "legacy binary Office documents are preserved raw but are not interpreted",
        )
        return _to_result(content, detected, declared_type, filename, output, active_limits)
    if detected == ZIP:
        output = _failure(
            "unsupported",
            "generic_archive",
            "generic archives require an explicit import workflow and were not expanded",
            metadata=zip_metadata,
        )
        return _to_result(content, detected, declared_type, filename, output, active_limits)

    try:
        if detected in _TEXT_TYPES or detected.startswith("text/"):
            output = _process_textual(content, detected, declared_type, active_limits)
        elif detected == PDF:
            output = _process_pdf(content, active_limits)
        elif detected == DOCX:
            output = _process_docx(content, active_limits)
        elif detected == XLSX:
            output = _process_xlsx(content, active_limits)
        elif detected == PPTX:
            output = _process_pptx(content, active_limits)
        elif detected in _IMAGE_TYPES:
            output = _process_image(content, active_limits)
        elif detected == "image/svg+xml":
            output = _failure(
                "unsupported",
                "active_vector_image",
                "SVG is preserved raw but is not rendered or served inline",
            )
        else:
            output = _failure(
                "unsupported",
                "unsupported_media_type",
                f"no safe parser is configured for {detected}",
            )
    except _DependencyUnavailable as exc:
        output = _failure("unsupported", "parser_unavailable", str(exc))
    except Exception as exc:
        output = _failure("corrupt", "parser_error", _exception_detail(exc, active_limits))

    if zip_metadata and output.metadata is not None:
        output.metadata = {**zip_metadata, **output.metadata}
    elif zip_metadata:
        output.metadata = zip_metadata
    return _to_result(content, detected, declared_type, filename, output, active_limits)


def _payload_bytes(payload: bytes | str | Mapping[str, Any]) -> bytes:
    if isinstance(payload, bytes):
        return payload
    if isinstance(payload, str):
        return payload.encode("utf-8")
    if isinstance(payload, Mapping):
        return json.dumps(
            payload,
            ensure_ascii=False,
            sort_keys=True,
            separators=(",", ":"),
            default=_safe_scalar,
        ).encode("utf-8")
    raise TypeError("payload must be bytes, str, or a mapping")


def _version_key(pipeline_version: str) -> str:
    version = pipeline_version.strip()
    if not version:
        raise ValueError("pipeline_version must not be empty")
    readable = re.sub(r"[^A-Za-z0-9._-]+", "-", version).strip("-.")[:40] or "pipeline"
    digest = hashlib.sha256(version.encode("utf-8")).hexdigest()[:12]
    return f"{readable}-{digest}"


def _safe_output_name(filename: str, compress: bool) -> str:
    if not _SAFE_OUTPUT_NAME_RE.fullmatch(filename):
        raise ValueError("filename must be a simple relative file name")
    return f"{filename}.gz" if compress and not filename.endswith(".gz") else filename


def _contained_path(root: Path, relative_path: str) -> Path:
    relative = Path(relative_path)
    if relative.is_absolute():
        raise ValueError("processed output path must be relative")
    resolved_root = root.expanduser().resolve()
    candidate = (resolved_root / relative).resolve()
    if not candidate.is_relative_to(resolved_root):
        raise ValueError("processed output path escapes the configured root")
    return candidate


def write_processed_output(
    processed_root: str | Path,
    content_hash: str,
    pipeline_version: str,
    payload: bytes | str | Mapping[str, Any],
    *,
    filename: str = "manifest.json",
    compress: bool = True,
) -> str:
    """Atomically write a versioned derived artifact and return its relative path."""

    if not _HEX_64.fullmatch(content_hash):
        raise ValueError("content_hash must be a 64-character hexadecimal SHA-256")
    digest = content_hash.lower()
    output_name = _safe_output_name(filename, compress)
    relative = (
        Path(digest[:2]) / digest[2:4] / digest / _version_key(pipeline_version) / output_name
    )
    root = Path(processed_root)
    target = _contained_path(root, relative.as_posix())
    target.parent.mkdir(parents=True, exist_ok=True)

    original = _payload_bytes(payload)
    stored = gzip.compress(original, mtime=0) if compress else original
    descriptor, temporary_name = tempfile.mkstemp(
        dir=target.parent, prefix=f".{target.name}.", suffix=".tmp"
    )
    temporary = Path(temporary_name)
    try:
        with os.fdopen(descriptor, "wb") as stream:
            stream.write(stored)
            stream.flush()
            os.fsync(stream.fileno())
        os.replace(temporary, target)
    except Exception:
        temporary.unlink(missing_ok=True)
        raise
    return relative.as_posix()


def read_processed_output(
    processed_root: str | Path,
    relative_path: str,
    *,
    max_bytes: int = 50 * 1024 * 1024,
) -> bytes:
    """Read a derived artifact with containment and decompression-size checks."""

    if max_bytes <= 0:
        raise ValueError("max_bytes must be positive")
    target = _contained_path(Path(processed_root), relative_path)
    if not target.is_file():
        raise FileNotFoundError(relative_path)
    if target.suffix != ".gz":
        if target.stat().st_size > max_bytes:
            raise ValueError("processed output exceeds the configured read limit")
        return target.read_bytes()
    with gzip.open(target, "rb") as stream:
        content = stream.read(max_bytes + 1)
    if len(content) > max_bytes:
        raise ValueError("decompressed processed output exceeds the configured read limit")
    return content
